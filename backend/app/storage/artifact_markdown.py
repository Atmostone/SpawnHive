"""Universal artifact → Markdown converter (SPA-71).

Turns a result deliverable (docx/pdf/pptx/xlsx/csv/ics/json) into Markdown so the
LLM judge can read it (instead of ``(binary file, content not shown)``) and a
human can read it rendered in the UI.

Pure/synchronous (CPU-bound); callers run it via ``asyncio.to_thread``. Every
converter is best-effort and **never raises** — on any failure (corrupt file,
missing lib, unsupported type) it degrades to ``None`` and the caller emits its
own note. Per-type libraries (no MarkItDown): pypdf, python-docx, python-pptx,
openpyxl, icalendar; csv/json via stdlib.
"""

import csv as _csv
import io
import json
import logging
import os

logger = logging.getLogger(__name__)

# Extensions we know how to convert. Anything else (images, archives, audio,
# unknown binary) → to_markdown returns None → the caller emits the binary note.
# Plain text (.md/.txt/.py/…) is intentionally NOT here: the judge already reads
# those cheaply via read_result_file_text, and the review UI falls back to that
# text — so there's no reason to force a full re-read through the converter.
_DOC_EXTS = frozenset({".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".ics", ".json"})

# Hard cap on converter OUTPUT length so a pathological 100k-row sheet can't blow
# the prompt before the judge-layer caps even apply (the judge truncates again).
_MAX_MARKDOWN_CHARS = 200_000

# Per-table bounds for spreadsheet/csv → keep Markdown (and memory) bounded.
_MAX_TABLE_ROWS = 200
_MAX_TABLE_COLS = 50

# Mirror of minio_client._MAX_CONVERT_BYTES (kept local to avoid import coupling).
_MAX_CONVERT_BYTES = 25 * 1024 * 1024


def _ext(name: str) -> str:
    return os.path.splitext(name)[1].lower()


def is_convertible(name: str) -> bool:
    """True if the extension is one we know how to convert.

    Drives the judge's routing: even a file the partial text reader decoded
    (e.g. .csv/.json — no NUL byte) should still go through the converter to get
    a clean GFM table / pretty JSON from a *full* read."""
    return _ext(name) in _DOC_EXTS


# --- GFM table helper (shared by docx tables / pptx tables / xlsx / csv) -------


def _esc(cell: object) -> str:
    s = "" if cell is None else str(cell)
    return s.replace("\\", "\\\\").replace("|", "\\|").replace("\n", " ").replace("\r", " ")


def _rows_to_gfm(rows: list[list]) -> str:
    """Render rows (first row = header) as a GFM table, bounded by row/col caps."""
    rows = [r for r in rows[:_MAX_TABLE_ROWS]]
    ncols = min(_MAX_TABLE_COLS, max((len(r) for r in rows), default=0))
    if ncols == 0:
        return ""

    def fmt(r: list) -> str:
        cells = [_esc(r[i]) if i < len(r) else "" for i in range(ncols)]
        return "| " + " | ".join(cells) + " |"

    lines = [fmt(rows[0]), "| " + " | ".join(["---"] * ncols) + " |"]
    lines += [fmt(r) for r in rows[1:]]
    return "\n".join(lines)


# --- per-type converters (each may raise; to_markdown wraps them) --------------


def _pdf_to_md(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _docx_to_md(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name if para.style else "") or ""
        if style.lower().startswith("heading"):
            digits = "".join(ch for ch in style if ch.isdigit())
            level = min(int(digits), 6) if digits.isdigit() else 1
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        gfm = _rows_to_gfm([[cell.text for cell in row.cells] for row in table.rows])
        if gfm:
            parts.append(gfm)
    return "\n\n".join(parts)


def _pptx_to_md(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                gfm = _rows_to_gfm(
                    [[cell.text for cell in row.cells] for row in shape.table.rows]
                )
                if gfm:
                    parts.append(gfm)
    return "\n\n".join(parts)


def _xlsx_to_md(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            parts.append(f"## {ws.title}")
            rows: list[list] = []
            truncated = False
            for r in ws.iter_rows(values_only=True):
                rows.append(list(r))
                if len(rows) >= _MAX_TABLE_ROWS:
                    truncated = True
                    break
            gfm = _rows_to_gfm(rows)
            if gfm:
                parts.append(gfm)
                if truncated:
                    parts.append(f"_… rows truncated at {_MAX_TABLE_ROWS}_")
    finally:
        wb.close()
    return "\n\n".join(parts)


def _csv_to_md(data: bytes) -> str:
    reader = _csv.reader(io.StringIO(data.decode("utf-8", errors="replace")))
    rows: list[list] = []
    truncated = False
    for r in reader:
        rows.append(r)
        if len(rows) >= _MAX_TABLE_ROWS:
            truncated = True
            break
    gfm = _rows_to_gfm(rows)
    if gfm and truncated:
        gfm += f"\n\n_… rows truncated at {_MAX_TABLE_ROWS}_"
    return gfm


def _ics_to_md(data: bytes) -> str:
    from icalendar import Calendar

    cal = Calendar.from_ical(data)
    events: list[str] = []
    for comp in cal.walk("VEVENT"):
        summary = str(comp.get("summary", "")).strip()
        lines = [f"### {summary or '(untitled event)'}"]
        for label, key in (("Start", "dtstart"), ("End", "dtend")):
            val = comp.get(key)
            if val is not None:
                lines.append(f"- **{label}:** {getattr(val, 'dt', val)}")
        for label, key in (("Location", "location"), ("Description", "description")):
            val = str(comp.get(key, "")).strip()
            if val:
                lines.append(f"- **{label}:** {val}")
        events.append("\n".join(lines))
    return "\n\n".join(events)


def _json_to_md(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    try:
        obj = json.loads(text)
    except Exception:
        return text  # not valid JSON — show as-is
    return "```json\n" + json.dumps(obj, indent=2, ensure_ascii=False) + "\n```"


_CONVERTERS = {
    ".pdf": _pdf_to_md,
    ".docx": _docx_to_md,
    ".pptx": _pptx_to_md,
    ".xlsx": _xlsx_to_md,
    ".csv": _csv_to_md,
    ".ics": _ics_to_md,
    ".json": _json_to_md,
}


def to_markdown(name: str, data: bytes) -> str | None:
    """Convert artifact bytes to Markdown by file extension.

    Returns Markdown text, or ``None`` when the type is unconvertible (images,
    archives, unknown binary), the file is empty, conversion fails, or the output
    is blank. **Never raises.** Synchronous; call via ``asyncio.to_thread``."""
    converter = _CONVERTERS.get(_ext(name))
    if converter is None or not data:
        return None
    try:
        md = converter(data)
    except Exception as e:  # noqa: BLE001 — a converter must never break the judge
        logger.warning(f"artifact_markdown: failed to convert {name}: {e}")
        return None
    if not md or not md.strip():
        return None
    if len(md) > _MAX_MARKDOWN_CHARS:
        md = md[:_MAX_MARKDOWN_CHARS] + "\n\n…[truncated]"
    return md


def result_file_markdown(s3_path: str, max_bytes: int = _MAX_CONVERT_BYTES) -> str | None:
    """Read full bytes from MinIO and convert to Markdown (SPA-71).

    Returns ``None`` on oversize file, read error, unconvertible type, or
    conversion failure. Synchronous; call via ``asyncio.to_thread``.

    Future hook: cache the produced Markdown beside the file at
    ``results-md/<task_id>/<file>.md`` to skip re-conversion on re-judge — not
    built in v1, the judge runs ~once per task so the parse cost is dwarfed by
    the LLM call."""
    name = s3_path.split("/", 2)[-1]  # strip the results/<task_id>/ prefix
    if not is_convertible(name):
        return None
    from app.storage.minio_client import read_result_file_bytes

    try:
        data = read_result_file_bytes(s3_path, max_bytes)
    except Exception as e:  # noqa: BLE001 — storage errors must not break the eval
        logger.warning(f"artifact_markdown: could not read {s3_path}: {e}")
        return None
    if not data:
        return None
    return to_markdown(name, data)
