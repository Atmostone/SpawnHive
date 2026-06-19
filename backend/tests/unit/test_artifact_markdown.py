"""Unit tests for the artifact → Markdown converter (SPA-71) — pure, no MinIO.

Builds tiny in-memory fixtures per type and asserts the converter produces
readable Markdown, degrades to None on corrupt/unknown/empty input, and never
raises."""

from __future__ import annotations

import io

from app.storage import artifact_markdown as am


def test_is_convertible_routing():
    assert am.is_convertible("report.docx")
    assert am.is_convertible("data/sheet.xlsx")
    assert am.is_convertible("a/b/c.csv")
    assert not am.is_convertible("notes.md")  # plain text → judge reads directly
    assert not am.is_convertible("logo.png")
    assert not am.is_convertible("archive.zip")


def test_csv_to_gfm_table():
    md = am.to_markdown("data.csv", b"name,score\nAlice,9\nBob,7\n")
    assert md is not None
    assert "| name | score |" in md
    assert "| --- | --- |" in md
    assert "| Alice | 9 |" in md


def test_json_pretty_fenced():
    md = am.to_markdown("payload.json", b'{"k": 1, "list": [1, 2]}')
    assert md is not None
    assert "```json" in md
    assert '"k": 1' in md


def test_json_invalid_falls_back_to_raw_text():
    md = am.to_markdown("payload.json", b"{not valid json")
    assert md is not None
    assert "not valid json" in md


def test_ics_lists_events():
    ics = (
        b"BEGIN:VCALENDAR\r\nVERSION:2.0\r\nPRODID:-//test//EN\r\n"
        b"BEGIN:VEVENT\r\nSUMMARY:Team Standup\r\nLOCATION:Room 4\r\n"
        b"DTSTART:20260620T090000Z\r\nEND:VEVENT\r\nEND:VCALENDAR\r\n"
    )
    md = am.to_markdown("cal.ics", ics)
    assert md is not None
    assert "Team Standup" in md
    assert "Room 4" in md


def test_docx_paragraphs_and_tables():
    from docx import Document

    doc = Document()
    doc.add_heading("Quarterly Report", level=1)
    doc.add_paragraph("The Q3 launch moves to Friday.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "100"
    buf = io.BytesIO()
    doc.save(buf)

    md = am.to_markdown("report.docx", buf.getvalue())
    assert md is not None
    assert "Quarterly Report" in md
    assert "The Q3 launch moves to Friday." in md
    assert "| Metric | Value |" in md
    assert "| Revenue | 100 |" in md


def test_xlsx_per_sheet_tables():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["region", "total"])
    ws.append(["EU", 42])
    buf = io.BytesIO()
    wb.save(buf)

    md = am.to_markdown("book.xlsx", buf.getvalue())
    assert md is not None
    assert "## Sales" in md
    assert "| region | total |" in md
    assert "| EU | 42 |" in md


def test_pptx_slides_as_sections():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Launch Plan"
    buf = io.BytesIO()
    prs.save(buf)

    md = am.to_markdown("deck.pptx", buf.getvalue())
    assert md is not None
    assert "## Slide 1" in md
    assert "Launch Plan" in md


def test_corrupt_docx_returns_none():
    assert am.to_markdown("report.docx", b"this is not a real docx file") is None


def test_unknown_extension_returns_none():
    assert am.to_markdown("logo.png", b"\x89PNG\r\n\x1a\n") is None


def test_empty_data_returns_none():
    assert am.to_markdown("data.csv", b"") is None


def test_output_truncated_to_cap():
    big = "row,val\n" + "\n".join(f"k,{i}" for i in range(500))
    md = am.to_markdown("big.csv", big.encode())
    assert md is not None
    # capped at _MAX_TABLE_ROWS rows → truncation note present
    assert "rows truncated" in md


def test_gfm_escapes_pipes_and_newlines():
    md = am.to_markdown("data.csv", b'a,b\n"x|y","line1\nline2"\n')
    assert md is not None
    assert "\\|" in md  # pipe escaped so it doesn't break the table
    assert "line1 line2" in md  # embedded newline flattened to a space
